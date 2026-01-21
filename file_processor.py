"""
File Processor for handling various submission formats.
Supports: .py, .ipynb, .zip (containing ipynb/py/pptx), .pptx
"""
import os
import json
import zipfile
import tempfile
import shutil
from pathlib import Path
from typing import Dict, Any, List
from dataclasses import dataclass


@dataclass
class SubmissionContent:
    """Represents extracted content from a submission."""
    filename: str
    code: str = ""
    markdown: str = ""
    outputs: str = ""
    presentation: str = ""
    file_type: str = ""
    errors: List[str] = None

    def __post_init__(self):
        if self.errors is None:
            self.errors = []

    def get_full_content(self) -> str:
        """Get all content combined for evaluation."""
        parts = []
        if self.code:
            parts.append(f"=== CODE ===\n{self.code}")
        if self.markdown:
            parts.append(f"=== MARKDOWN/DOCUMENTATION ===\n{self.markdown}")
        if self.outputs:
            parts.append(f"=== OUTPUTS/RESULTS ===\n{self.outputs}")
        if self.presentation:
            parts.append(f"=== PRESENTATION ===\n{self.presentation}")
        return "\n\n".join(parts)


class FileProcessor:
    """Process various file formats for grading."""

    def __init__(self):
        """Initialize the file processor."""
        self.temp_dirs = []

    def __del__(self):
        """Cleanup temporary directories."""
        for temp_dir in self.temp_dirs:
            try:
                shutil.rmtree(temp_dir)
            except Exception:
                pass

    def process_file(self, file_path: str) -> SubmissionContent:
        """
        Process a submission file and extract content.

        Args:
            file_path: Path to the submission file.

        Returns:
            SubmissionContent with extracted data.
        """
        path = Path(file_path)
        filename = path.name
        extension = path.suffix.lower()

        if extension == ".py":
            return self._process_python(file_path, filename)
        elif extension == ".ipynb":
            return self._process_notebook(file_path, filename)
        elif extension == ".zip":
            return self._process_zip(file_path, filename)
        elif extension == ".pptx":
            return self._process_pptx(file_path, filename)
        else:
            return SubmissionContent(
                filename=filename,
                errors=[f"Unsupported file type: {extension}"]
            )

    def process_uploaded_file(self, uploaded_file) -> SubmissionContent:
        """
        Process an uploaded file from Streamlit.

        Args:
            uploaded_file: Streamlit UploadedFile object.

        Returns:
            SubmissionContent with extracted data.
        """
        # Save to temporary file
        temp_dir = tempfile.mkdtemp()
        self.temp_dirs.append(temp_dir)

        temp_path = os.path.join(temp_dir, uploaded_file.name)
        with open(temp_path, 'wb') as f:
            f.write(uploaded_file.getvalue())

        return self.process_file(temp_path)

    def _process_python(self, file_path: str, filename: str) -> SubmissionContent:
        """Process a Python file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                code = f.read()

            # Extract docstrings and comments as documentation
            markdown = self._extract_python_docs(code)

            return SubmissionContent(
                filename=filename,
                code=code,
                markdown=markdown,
                file_type="python"
            )
        except Exception as e:
            return SubmissionContent(
                filename=filename,
                errors=[f"Error reading Python file: {str(e)}"]
            )

    def _extract_python_docs(self, code: str) -> str:
        """Extract docstrings and comments from Python code."""
        import re

        docs = []

        # Extract module docstring
        module_doc = re.match(r'^["\'][\'"]{2}(.*?)["\'][\'"]{2}', code, re.DOTALL)
        if module_doc:
            docs.append(f"Module Documentation:\n{module_doc.group(1)}")

        # Extract function/class docstrings
        docstrings = re.findall(r'def \w+.*?:\s*["\'][\'"]{2}(.*?)["\'][\'"]{2}', code, re.DOTALL)
        docstrings += re.findall(r'class \w+.*?:\s*["\'][\'"]{2}(.*?)["\'][\'"]{2}', code, re.DOTALL)

        for doc in docstrings:
            docs.append(doc.strip())

        # Extract comments
        comments = re.findall(r'#\s*(.+)$', code, re.MULTILINE)
        if comments:
            docs.append("Comments:\n" + "\n".join(comments[:20]))  # Limit comments

        return "\n\n".join(docs)

    def _process_notebook(self, file_path: str, filename: str) -> SubmissionContent:
        """Process a Jupyter notebook."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                notebook = json.load(f)

            code_cells = []
            markdown_cells = []
            outputs = []

            cells = notebook.get('cells', [])

            for i, cell in enumerate(cells):
                cell_type = cell.get('cell_type', '')
                source = ''.join(cell.get('source', []))

                if cell_type == 'code':
                    code_cells.append(f"# Cell {i + 1}\n{source}")

                    # Extract outputs
                    cell_outputs = cell.get('outputs', [])
                    for output in cell_outputs:
                        if 'text' in output:
                            text = ''.join(output['text'])
                            outputs.append(f"[Cell {i + 1} Output]:\n{text}")
                        elif 'data' in output:
                            data = output['data']
                            if 'text/plain' in data:
                                text = ''.join(data['text/plain'])
                                outputs.append(f"[Cell {i + 1} Output]:\n{text}")
                            if 'text/html' in data:
                                outputs.append(f"[Cell {i + 1}]: HTML output (table/visualization)")
                            if 'image/png' in data:
                                outputs.append(f"[Cell {i + 1}]: Image output (plot/chart)")

                elif cell_type == 'markdown':
                    markdown_cells.append(source)

            return SubmissionContent(
                filename=filename,
                code="\n\n".join(code_cells),
                markdown="\n\n".join(markdown_cells),
                outputs="\n\n".join(outputs),
                file_type="notebook"
            )
        except Exception as e:
            return SubmissionContent(
                filename=filename,
                errors=[f"Error reading notebook: {str(e)}"]
            )

    def _process_zip(self, file_path: str, filename: str) -> SubmissionContent:
        """Process a ZIP file containing submissions."""
        temp_dir = tempfile.mkdtemp()
        self.temp_dirs.append(temp_dir)

        try:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)

            # Find all relevant files
            all_code = []
            all_markdown = []
            all_outputs = []
            all_presentations = []
            processed_files = []

            for root, dirs, files in os.walk(temp_dir):
                # Skip hidden directories and __MACOSX
                dirs[:] = [d for d in dirs if not d.startswith('.') and d != '__MACOSX']

                for file in files:
                    if file.startswith('.'):
                        continue

                    file_path_inner = os.path.join(root, file)
                    ext = Path(file).suffix.lower()

                    if ext == '.py':
                        content = self._process_python(file_path_inner, file)
                        if content.code:
                            all_code.append(f"# File: {file}\n{content.code}")
                        if content.markdown:
                            all_markdown.append(content.markdown)
                        processed_files.append(file)

                    elif ext == '.ipynb':
                        content = self._process_notebook(file_path_inner, file)
                        if content.code:
                            all_code.append(f"# Notebook: {file}\n{content.code}")
                        if content.markdown:
                            all_markdown.append(f"# From {file}:\n{content.markdown}")
                        if content.outputs:
                            all_outputs.append(f"# From {file}:\n{content.outputs}")
                        processed_files.append(file)

                    elif ext == '.pptx':
                        content = self._process_pptx(file_path_inner, file)
                        if content.presentation:
                            all_presentations.append(content.presentation)
                        processed_files.append(file)

            return SubmissionContent(
                filename=filename,
                code="\n\n".join(all_code),
                markdown="\n\n".join(all_markdown),
                outputs="\n\n".join(all_outputs),
                presentation="\n\n".join(all_presentations),
                file_type=f"zip (contains: {', '.join(processed_files)})"
            )

        except Exception as e:
            return SubmissionContent(
                filename=filename,
                errors=[f"Error extracting ZIP: {str(e)}"]
            )

    def _process_pptx(self, file_path: str, filename: str) -> SubmissionContent:
        """Process a PowerPoint file."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slides_content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"=== Slide {slide_num} ==="]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # Check for tables
                    if shape.has_table:
                        table = shape.table
                        table_text = ["[Table]"]
                        for row in table.rows:
                            row_text = [cell.text for cell in row.cells]
                            table_text.append(" | ".join(row_text))
                        slide_text.append("\n".join(table_text))

                slides_content.append("\n".join(slide_text))

            return SubmissionContent(
                filename=filename,
                presentation="\n\n".join(slides_content),
                file_type="powerpoint"
            )

        except ImportError:
            return SubmissionContent(
                filename=filename,
                errors=["python-pptx not installed. Run: pip install python-pptx"]
            )
        except Exception as e:
            return SubmissionContent(
                filename=filename,
                errors=[f"Error reading PowerPoint: {str(e)}"]
            )

    def process_dataset(self, file_path: str) -> Dict[str, Any]:
        """
        Process a dataset file and extract summary information.

        Args:
            file_path: Path to the dataset file.

        Returns:
            Dict with dataset information.
        """
        import pandas as pd

        path = Path(file_path)
        ext = path.suffix.lower()

        try:
            if ext == '.csv':
                df = pd.read_csv(file_path)
            elif ext in ['.xlsx', '.xls']:
                df = pd.read_excel(file_path)
            elif ext == '.json':
                df = pd.read_json(file_path)
            else:
                return {
                    "success": False,
                    "error": f"Unsupported dataset format: {ext}"
                }

            # Generate summary
            summary = {
                "success": True,
                "filename": path.name,
                "shape": df.shape,
                "columns": list(df.columns),
                "dtypes": df.dtypes.to_dict(),
                "head": df.head(5).to_string(),
                "describe": df.describe().to_string(),
                "missing": df.isnull().sum().to_dict()
            }

            # Create a text description
            summary["description"] = f"""
Dataset: {path.name}
Rows: {df.shape[0]}, Columns: {df.shape[1]}

Columns: {', '.join(df.columns)}

Data Types:
{chr(10).join([f'  - {col}: {dtype}' for col, dtype in df.dtypes.items()])}

Sample Data (first 5 rows):
{df.head(5).to_string()}

Statistical Summary:
{df.describe().to_string()}

Missing Values:
{chr(10).join([f'  - {col}: {count}' for col, count in df.isnull().sum().items() if count > 0]) or '  None'}
"""
            return summary

        except Exception as e:
            return {
                "success": False,
                "error": f"Error reading dataset: {str(e)}"
            }

    def process_uploaded_dataset(self, uploaded_file) -> Dict[str, Any]:
        """Process an uploaded dataset file from Streamlit."""
        temp_dir = tempfile.mkdtemp()
        self.temp_dirs.append(temp_dir)

        temp_path = os.path.join(temp_dir, uploaded_file.name)
        with open(temp_path, 'wb') as f:
            f.write(uploaded_file.getvalue())

        return self.process_dataset(temp_path)
